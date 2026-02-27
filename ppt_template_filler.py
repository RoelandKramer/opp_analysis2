# ============================================================
# file: ppt_template_filler.py
# ============================================================
from __future__ import annotations

import io
import os
import re
from dataclasses import dataclass
from typing import Dict, Iterable, List, Optional, Tuple

import matplotlib.figure
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


@dataclass(frozen=True)
class CropSpec:
    left: float = 0.0
    right: float = 0.0
    top: float = 0.0
    bottom: float = 0.0


CROP_BY_SHAPE_NAME: Dict[str, CropSpec] = {
    "PH_Corners_left_positions_vis": CropSpec(left=0.0162, top=0.1850, bottom=0.0230),
    "PH_Corners_right_positions_vis": CropSpec(left=0.0162, right=0.0700, top=0.1287, bottom=0.2993),
    "PH_Corners_left_shots_vis": CropSpec(bottom=0.2514),
    "PH_Corners_right_shots_vis": CropSpec(bottom=0.2514),
    "PH_def_left": CropSpec(bottom=0.2514),
    "PH_def_right": CropSpec(bottom=0.2514),
    # add cropping here later if needed:
    # "PH_att_c_headers": CropSpec(...),
    # "PH_def_c_headers": CropSpec(...),
}


# If shape-name placement fails, we try token placement as a fallback.
FALLBACK_SHAPE_TO_TOKEN: Dict[str, str] = {
    "PH_att_c_headers": "{att_corners_headers}",
    "PH_def_c_headers": "{def_corners_headers}",
}


def _hex_to_rgb(hex_color: str) -> RGBColor:
    s = (hex_color or "").strip().lstrip("#")
    if len(s) != 6:
        return RGBColor(255, 255, 255)
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def fig_to_png_bytes_and_close(fig, *, dpi=240) -> bytes:
    try:
        return fig_to_png_bytes(fig, dpi=dpi)
    finally:
        try:
            plt.close(fig)
        except Exception:
            pass

def fig_to_png_bytes_labels_and_close(fig, *, dpi=240) -> bytes:
    try:
        return fig_to_png_bytes_labels(fig, dpi=dpi)
    finally:
        try:
            plt.close(fig)
        except Exception:
            pass
            
def fig_to_png_bytes(fig: matplotlib.figure.Figure, *, dpi: int = 240) -> bytes:
    try:
        fig.subplots_adjust(left=0, right=1, top=1, bottom=0, wspace=0, hspace=0)
        fig.patch.set_alpha(0)
        for ax in fig.axes:
            ax.set_position([0, 0, 1, 1])
            ax.margins(0)
    except Exception:
        pass

    bio = io.BytesIO()
    fig.savefig(bio, format="png", dpi=dpi, transparent=True, bbox_inches=None, pad_inches=0)
    return bio.getvalue()


def fig_to_png_bytes_labels(fig: matplotlib.figure.Figure, *, dpi: int = 240) -> bytes:
    bio = io.BytesIO()
    fig.savefig(bio, format="png", dpi=dpi, transparent=True, bbox_inches="tight", pad_inches=0.08)
    return bio.getvalue()


def _clamp(v: float, lo: float, hi: float) -> float:
    return max(lo, min(hi, v))


def crop_png_bytes(img_bytes: bytes, *, crop: CropSpec) -> bytes:
    im = Image.open(io.BytesIO(img_bytes)).convert("RGBA")
    w, h = im.size

    l = int(round(_clamp(crop.left, 0.0, 0.45) * w))
    r = int(round(_clamp(crop.right, 0.0, 0.45) * w))
    t = int(round(_clamp(crop.top, 0.0, 0.45) * h))
    b = int(round(_clamp(crop.bottom, 0.0, 0.45) * h))

    x0 = l
    y0 = t
    x1 = max(x0 + 1, w - r)
    y1 = max(y0 + 1, h - b)

    im = im.crop((x0, y0, x1, y1))

    out = io.BytesIO()
    im.save(out, format="PNG")
    return out.getvalue()


ShapeMapped = Tuple[object, int, int, int, int, Optional[object]]


def _children_bbox(group_shape) -> Tuple[int, int, int, int]:
    min_l, min_t = 10**18, 10**18
    max_r, max_b = -10**18, -10**18

    for ch in group_shape.shapes:
        l, t = int(ch.left), int(ch.top)
        r, b = l + int(ch.width), t + int(ch.height)
        min_l = min(min_l, l)
        min_t = min(min_t, t)
        max_r = max(max_r, r)
        max_b = max(max_b, b)

    if min_l == 10**18:
        return 0, 0, 1, 1
    return min_l, min_t, max_r, max_b


def iter_shapes_mapped(
    container,
    *,
    base_left: int = 0,
    base_top: int = 0,
    base_sx: float = 1.0,
    base_sy: float = 1.0,
    parent_group: Optional[object] = None,
) -> Iterable[ShapeMapped]:
    for shp in container.shapes:
        rel_l, rel_t = int(shp.left), int(shp.top)
        rel_w, rel_h = int(shp.width), int(shp.height)

        abs_left = int(round(base_left + rel_l * base_sx))
        abs_top = int(round(base_top + rel_t * base_sy))
        abs_w = int(round(rel_w * base_sx))
        abs_h = int(round(rel_h * base_sy))

        yield shp, abs_left, abs_top, abs_w, abs_h, parent_group

        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            min_l, min_t, max_r, max_b = _children_bbox(shp)
            internal_w = max(1, max_r - min_l)
            internal_h = max(1, max_b - min_t)

            sx = abs_w / internal_w
            sy = abs_h / internal_h

            child_base_left = abs_left - int(round(min_l * sx))
            child_base_top = abs_top - int(round(min_t * sy))

            yield from iter_shapes_mapped(
                shp,
                base_left=child_base_left,
                base_top=child_base_top,
                base_sx=sx,
                base_sy=sy,
                parent_group=shp,
            )


def _shape_text(shp) -> str:
    if not getattr(shp, "has_text_frame", False):
        return ""
    try:
        return shp.text_frame.text or ""
    except Exception:
        return ""


def _replace_text_in_shape(shp, replacements: Dict[str, str]) -> None:
    if not getattr(shp, "has_text_frame", False):
        return

    tf = shp.text_frame

    for para in tf.paragraphs:
        for run in para.runs:
            txt = run.text or ""
            for k, v in replacements.items():
                if k in txt:
                    txt = txt.replace(k, v)
            run.text = txt

    for para in tf.paragraphs:
        whole = "".join(run.text for run in para.runs)
        replaced = whole
        for k, v in replacements.items():
            replaced = replaced.replace(k, v)
        if replaced != whole:
            for run in list(para.runs):
                run.text = ""
            if para.runs:
                para.runs[0].text = replaced
            else:
                para.add_run().text = replaced


def _clear_token_text(shp, token: str) -> None:
    if not getattr(shp, "has_text_frame", False):
        return
    tf = shp.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            if run.text and token in run.text:
                run.text = run.text.replace(token, "")


def _delete_shape(shp) -> None:
    try:
        el = shp.element
        el.getparent().remove(el)
    except Exception:
        pass


def _set_slide_background(slide, *, bg_hex: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(bg_hex)


def _set_shape_fill(shp, *, hex_color: str) -> None:
    try:
        fill = shp.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(hex_color)
    except Exception:
        pass


def _find_shapes_with_token(slide, token: str) -> List[ShapeMapped]:
    out: List[ShapeMapped] = []
    for shp, ax, ay, aw, ah, parent in iter_shapes_mapped(slide):
        if token in _shape_text(shp):
            out.append((shp, ax, ay, aw, ah, parent))
    return out


def _shape_text_stripped(shp) -> str:
    return (_shape_text(shp) or "").strip()


def _find_shapes_with_token_exact(slide, token: str) -> List[ShapeMapped]:
    out: List[ShapeMapped] = []
    for shp, ax, ay, aw, ah, parent in iter_shapes_mapped(slide):
        if _shape_text_stripped(shp) == token:
            out.append((shp, ax, ay, aw, ah, parent))
    return out


def _color_background_rectangles(slide, *, bg_hex: str, slide_w: int, slide_h: int) -> None:
    for shp, ax, ay, aw, ah, _ in iter_shapes_mapped(slide):
        if (
            ax <= int(slide_w * 0.05)
            and ay <= int(slide_h * 0.05)
            and aw >= int(slide_w * 0.90)
            and ah >= int(slide_h * 0.90)
        ):
            _set_shape_fill(shp, hex_color=bg_hex)


def _is_bar_candidate(ax: int, ay: int, aw: int, ah: int, slide_w: int, slide_h: int) -> bool:
    if aw < int(slide_w * 0.80):
        return False
    if ah > int(slide_h * 0.20):
        return False
    near_top = ay <= int(slide_h * 0.14)
    near_bottom = (ay + ah) >= int(slide_h * 0.86)
    return (near_top or near_bottom) and ax <= int(slide_w * 0.14)


def _color_and_clear_token_shapes(slide, token: str, *, hex_color: str) -> None:
    for shp, *_ in _find_shapes_with_token(slide, token):
        _set_shape_fill(shp, hex_color=hex_color)
        _clear_token_text(shp, token)


def _color_top_and_bottom_bars(slide, *, bar_hex: str, slide_w: int, slide_h: int) -> None:
    for shp, ax, ay, aw, ah, _ in iter_shapes_mapped(slide):
        if _is_bar_candidate(ax, ay, aw, ah, slide_w, slide_h):
            _set_shape_fill(shp, hex_color=bar_hex)

    _color_and_clear_token_shapes(slide, "{top_bar}", hex_color=bar_hex)
    _color_and_clear_token_shapes(slide, "{middle_bar}", hex_color=bar_hex)
    _color_and_clear_token_shapes(slide, "{bottom_bar}", hex_color=bar_hex)


Target = Tuple[object, int, int, int, int]


def _find_token_targets(slide, token: str, *, exact: bool = True) -> List[Target]:
    finder = _find_shapes_with_token_exact if exact else _find_shapes_with_token
    hits = finder(slide, token)

    out: List[Target] = []
    for token_shp, token_ax, token_ay, token_aw, token_ah, _ in hits:
        out.append((token_shp, token_ax, token_ay, token_aw, token_ah))
    return sorted(out, key=lambda t: (t[1], t[2]))


def _insert_picture_fill(slide, *, delete_shape: object, left: int, top: int, w: int, h: int, img_bytes: bytes) -> None:
    _delete_shape(delete_shape)
    slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width=w, height=h)


def _fill_token_with_images(slide, token: str, images: List[bytes], *, exact: bool = True) -> None:
    targets = _find_token_targets(slide, token, exact=exact)
    if not targets or not images:
        return

    for (del_shp, left, top, w, h), img in zip(targets, images):
        _insert_picture_fill(slide, delete_shape=del_shp, left=left, top=top, w=w, h=h, img_bytes=img)


def _normalize_name(s: Optional[str]) -> str:
    if not s:
        return ""
    # Normalize spaces and case
    s2 = re.sub(r"\s+", " ", s).strip().lower()
    return s2


def _find_shape_by_name_mapped(slide, shape_name: str) -> Optional[ShapeMapped]:
    target = _normalize_name(shape_name)
    for shp, ax, ay, aw, ah, parent in iter_shapes_mapped(slide):
        nm = _normalize_name(getattr(shp, "name", None))
        if nm == target:
            return shp, ax, ay, aw, ah, parent
    # second pass: allow substring match (covers weird suffix/prefix issues)
    for shp, ax, ay, aw, ah, parent in iter_shapes_mapped(slide):
        nm = _normalize_name(getattr(shp, "name", None))
        if target and nm and (target in nm or nm in target):
            return shp, ax, ay, aw, ah, parent
    return None


def _replace_named_shape_with_picture(slide, shape_name: str, img_bytes: bytes) -> bool:
    hit = _find_shape_by_name_mapped(slide, shape_name)
    if hit is None:
        return False

    shp, ax, ay, aw, ah, _ = hit
    _delete_shape(shp)
    slide.shapes.add_picture(io.BytesIO(img_bytes), ax, ay, width=aw, height=ah)
    return True


def _find_tables(slide) -> List:
    return [shp.table for shp in slide.shapes if getattr(shp, "has_table", False)]


def _write_df_to_ppt_table(table, df: pd.DataFrame) -> None:
    if df is None:
        df = pd.DataFrame()

    n_rows = len(table.rows)
    n_cols = len(table.columns)

    for r in range(1, n_rows):
        for c in range(n_cols):
            table.cell(r, c).text = ""

    if not df.empty:
        values = df.astype(str).replace({"nan": "-", "None": "-"}).values.tolist()
        max_write = max(0, n_rows - 1)
        for r in range(min(max_write, len(values))):
            row_vals = values[r]
            for c in range(min(n_cols, len(row_vals))):
                table.cell(r + 1, c).text = row_vals[c]

    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            if not cell.text_frame:
                continue
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)


TAKERS_LEFT_TABLE_SHAPE_NAME = "PH_takers_left"
TAKERS_RIGHT_TABLE_SHAPE_NAME = "PH_takers_right"


def _find_table_shape_by_name(prs: Presentation, shape_name: str):
    for slide in prs.slides:
        for shp in slide.shapes:
            if getattr(shp, "has_table", False) and getattr(shp, "name", None) == shape_name:
                return shp.table
    return None


def _fill_takers_tables(prs: Presentation, left_df: pd.DataFrame, right_df: pd.DataFrame) -> None:
    left_tbl = _find_table_shape_by_name(prs, TAKERS_LEFT_TABLE_SHAPE_NAME)
    right_tbl = _find_table_shape_by_name(prs, TAKERS_RIGHT_TABLE_SHAPE_NAME)

    if left_tbl is not None:
        _write_df_to_ppt_table(left_tbl, left_df)
    if right_tbl is not None:
        _write_df_to_ppt_table(right_tbl, right_df)

    if left_tbl is None and len(prs.slides) >= 1:
        t1 = _find_tables(prs.slides[0])
        if t1:
            _write_df_to_ppt_table(t1[0], left_df)

    if right_tbl is None and len(prs.slides) >= 2:
        t2 = _find_tables(prs.slides[1])
        if t2:
            _write_df_to_ppt_table(t2[0], right_df)


def _flatten_images_by_shape_name(images_by_shape_name: Optional[Dict[int, Dict[str, bytes]]]) -> Dict[str, bytes]:
    flat: Dict[str, bytes] = {}
    for _slide_idx, mapping in (images_by_shape_name or {}).items():
        for shape_name, img_bytes in (mapping or {}).items():
            flat[shape_name] = img_bytes
    return flat


@dataclass(frozen=True)
class FilledPptPayload:
    pptx_bytes: bytes
    filename: str


def fill_corner_template_pptx(
    *,
    template_pptx_path: str,
    team_name: str,
    team_primary_hex: str,
    team_secondary_hex: str,
    logo_path: Optional[str],
    meta_replacements: Dict[str, str],
    images_by_token: Dict[str, List[bytes]],
    images_by_shape_name: Optional[Dict[int, Dict[str, bytes]]] = None,
    left_takers_df: pd.DataFrame,
    right_takers_df: pd.DataFrame,
) -> FilledPptPayload:
    if not os.path.exists(template_pptx_path):
        raise FileNotFoundError(f"Template not found: {template_pptx_path}")

    prs = Presentation(template_pptx_path)
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    base_repl = dict(meta_replacements or {})
    base_repl.setdefault("{TEAM_NAME}", team_name)

    token_to_named = {
        "{Corners_left_positions_vis}": "PH_Corners_left_positions_vis",
        "{Corners_right_positions_vis}": "PH_Corners_right_positions_vis",
        "{Corners_left_shots_vis}": "PH_Corners_left_shots_vis",
        "{Corners_right_shots_vis}": "PH_Corners_right_shots_vis",
    }

    remaining_named = _flatten_images_by_shape_name(images_by_shape_name)
    remaining_named_original = dict(remaining_named)

    for slide in prs.slides:
        _set_slide_background(slide, bg_hex=team_secondary_hex)
        _color_background_rectangles(slide, bg_hex=team_secondary_hex, slide_w=slide_w, slide_h=slide_h)
        _color_top_and_bottom_bars(slide, bar_hex=team_primary_hex, slide_w=slide_w, slide_h=slide_h)

        for shp, *_ in iter_shapes_mapped(slide):
            _replace_text_in_shape(shp, base_repl)

        if logo_path and os.path.exists(logo_path):
            with open(logo_path, "rb") as f:
                logo_bytes = f.read()
            _fill_token_with_images(slide, "{LOGO}", [logo_bytes], exact=True)

        replaced_shape_names: set[str] = set()
        for shape_name, img_bytes in list(remaining_named.items()):
            hit = _find_shape_by_name_mapped(slide, shape_name)
            if hit is None:
                continue

            crop_spec = CROP_BY_SHAPE_NAME.get(shape_name)
            if crop_spec is not None:
                img_bytes = crop_png_bytes(img_bytes, crop=crop_spec)

            _replace_named_shape_with_picture(slide, shape_name, img_bytes)
            replaced_shape_names.add(shape_name)
            remaining_named.pop(shape_name, None)

        # token images
        for token, imgs in (images_by_token or {}).items():
            if token in ("{top_bar}", "{bottom_bar}", "{middle_bar}"):
                continue

            mapped_name = token_to_named.get(token)
            if mapped_name and mapped_name in replaced_shape_names:
                continue

            if token in ("{att_corners_headers}", "{def_corners_headers}"):
                _fill_token_with_images(slide, token, imgs, exact=False)
            else:
                _fill_token_with_images(slide, token, imgs, exact=True)

        for tok in ("{top_bar}", "{middle_bar}", "{bottom_bar}"):
            for shp, *_ in _find_shapes_with_token(slide, tok):
                _clear_token_text(shp, tok)

    # ✅ FINAL fallback: if PH_att/def_c_headers were NOT found by name, place them by token
    # (This covers master/layout shapes, or name mismatch, etc.)
    for shape_name, token in FALLBACK_SHAPE_TO_TOKEN.items():
        if shape_name in remaining_named:
            # try to place via token anywhere it exists (non-exact match)
            img_bytes = remaining_named[shape_name]
            for slide in prs.slides:
                _fill_token_with_images(slide, token, [img_bytes], exact=False)
            remaining_named.pop(shape_name, None)

    _fill_takers_tables(prs, left_takers_df, right_takers_df)

    out = io.BytesIO()
    prs.save(out)
    fname = f"OpponentAnalysis_Corners_{team_name.replace(' ', '_')}.pptx"
    return FilledPptPayload(pptx_bytes=out.getvalue(), filename=fname)
